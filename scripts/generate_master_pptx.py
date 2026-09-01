#!/usr/bin/env python3
"""
Master PowerPoint (.pptx) Presentation Generator for Chara
Generates a widescreen (16:9) academic light theme presentation with:
- Embedded HD Molecular Dynamics Video (md_simulation.mp4)
- 300 DPI Publication Figures (Fig 1 to Fig 6, RMSD Replicas, Rg Replicas)
- 3D Protein Structures (KRAS, c-MYC, PTPN11, TP53)
- Crisp data tables, benchmark comparisons, and clinical metrics
- Zero math bloat; structured around the 15 major biological & clinical steps
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color Palette (Academic Light)
BG_COLOR = RGBColor(248, 250, 252)       # #f8fafc Soft Slate
WHITE_COLOR = RGBColor(255, 255, 255)    # #ffffff Crisp Card
NAVY_COLOR = RGBColor(15, 23, 42)        # #0f172a Deep Navy
BLUE_COLOR = RGBColor(37, 99, 235)       # #2563eb Royal Blue
EMERALD_COLOR = RGBColor(22, 163, 74)    # #16a34a Emerald Green
CRIMSON_COLOR = RGBColor(220, 38, 38)    # #dc2626 Crimson Red
AMBER_COLOR = RGBColor(217, 119, 6)      # #d97706 Amber
BORDER_COLOR = RGBColor(226, 232, 240)   # #e2e8f0
GRAY_TEXT = RGBColor(100, 116, 139)      # #64748b

ROOT = Path(r"E:\Sharon")
LOGO_PATH = ROOT / "assets" / "iit-mandi-logo.png"
AUTHOR_PATH = ROOT / "assets" / "sharon-melhi.png"
PI_PATH = ROOT / "assets" / "dr-kharerin-hungyo.png"
VIDEO_PATH = ROOT / "md_simulation.mp4"

FIG1_PATH = ROOT / "Fig1_OOD_Performance.png"
FIG2_PATH = ROOT / "Fig2_KM_Survival.png"
FIG3_PATH = ROOT / "Fig3_Ablation_Impact.png"
FIG4_PATH = ROOT / "Fig4_Adversarial_Decay.png"
FIG5_PATH = ROOT / "Fig5_Dirichlet_Energy.png"
FIG6_PATH = ROOT / "Fig6_GSEA_Enrichment.png"
RMSD_PATH = ROOT / "RMSD_Replicas_300dpi.png"
RG_PATH = ROOT / "Rg_Replicas_300dpi.png"

KRAS_PATH = ROOT / "data" / "processed" / "images" / "KRAS_G12D_structure_PTM.png"
CMYC_PATH = ROOT / "data" / "processed" / "images" / "cMYC_MAX_structure_PTM.png"
PTPN11_PATH = ROOT / "data" / "processed" / "images" / "PTPN11_structure_PTM.png"
P53_PATH = ROOT / "data" / "processed" / "images" / "Mut_p53_structure_PTM.png"

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_COLOR
        bg_shape.line.fill.background() # No line

    def add_header(slide, title_text, category_text=""):
        # Header banner card
        header_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.95))
        header_card.fill.solid()
        header_card.fill.fore_color.rgb = WHITE_COLOR
        header_card.line.color.rgb = BORDER_COLOR
        header_card.line.width = Pt(1)

        # Logo
        if LOGO_PATH.exists():
            slide.shapes.add_picture(str(LOGO_PATH), Inches(0.95), Inches(0.52), height=Inches(0.7))

        # Title text
        tx_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.45), Inches(8.5), Inches(0.85))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = category_text.upper() if category_text else "IIT MANDI · CPG LAB"
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = BLUE_COLOR
        p1.font.name = "Arial"

        p2 = tf.add_paragraph()
        p2.text = title_text
        p2.font.size = Pt(17)
        p2.font.bold = True
        p2.font.color.rgb = NAVY_COLOR
        p2.font.name = "Arial"

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Main Center Card
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(0.8), Inches(10.333), Inches(5.9))
    card1.fill.solid()
    card1.fill.fore_color.rgb = WHITE_COLOR
    card1.line.color.rgb = BORDER_COLOR
    card1.line.width = Pt(1.5)

    if LOGO_PATH.exists():
        slide1.shapes.add_picture(str(LOGO_PATH), Inches(6.066), Inches(1.1), height=Inches(1.1))

    tx1 = slide1.shapes.add_textbox(Inches(1.8), Inches(2.3), Inches(9.733), Inches(2.2))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "COMPUTATIONAL & PHYSICAL GENOMICS LABORATORY · IIT MANDI"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BLUE_COLOR

    p = tf1.add_paragraph()
    p.text = "Chara: Molecular-Dynamics-Guided Survival Generalization"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY_COLOR

    p = tf1.add_paragraph()
    p.text = "Zero-Retraining Cross-Platform Cancer Prognosis via Biophysically Weighted Spectral Diffusion Networks"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY_TEXT

    # Authors Box
    tx_auth = slide1.shapes.add_textbox(Inches(2.5), Inches(4.8), Inches(8.333), Inches(1.5))
    tf_auth = tx_auth.text_frame
    tf_auth.word_wrap = True
    p = tf_auth.paragraphs[0]
    p.text = "Lead Author: Sharon Melhi Nadar    |    Principal Investigator: Dr. Kharerin Hungyo"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY_COLOR

    p = tf_auth.add_paragraph()
    p.text = "Indian Institute of Technology Mandi, Himachal Pradesh, India"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_TEXT

    # =========================================================================
    # SLIDE 2: THE GENERALIZATION CRISIS
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "The Cross-Platform Generalization Crisis", "The Clinical Problem")

    # Left Column Text
    card_l2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.3))
    card_l2.fill.solid()
    card_l2.fill.fore_color.rgb = WHITE_COLOR
    card_l2.line.color.rgb = BORDER_COLOR
    
    tx2_l = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf2_l = tx2_l.text_frame
    tf2_l.word_wrap = True
    
    p = tf2_l.paragraphs[0]
    p.text = "Why 99% of Published Oncology AI Models Fail:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = CRIMSON_COLOR

    items2 = [
        ("Catastrophic Domain Shift: ", "Models trained on Illumina RNA-seq (TCGA-LUAD, n=503) memorize platform-specific negative binomial probe distributions."),
        ("Microarray Failure: ", "When tested on historic Affymetrix Microarray biobanks (GSE31210, n=226), standard neural nets drop to C = 0.5537 and random forests to C = 0.4041."),
        ("The Biophysical Solution: ", "Sequencing probe chemistries shift every decade, but protein allosteric kinetics and binding thermodynamics never change. Chara uses physical simulations as an invariant inductive prior.")
    ]
    for bold_prefix, text in items2:
        p = tf2_l.add_paragraph()
        p.text = "• " + bold_prefix + text
        p.font.size = Pt(11.5)
        p.font.color.rgb = NAVY_COLOR
        p.space_before = Pt(8)

    # Right Column Figure 1
    card_r2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.6), Inches(5.8), Inches(5.3))
    card_r2.fill.solid()
    card_r2.fill.fore_color.rgb = WHITE_COLOR
    card_r2.line.color.rgb = BORDER_COLOR

    if FIG1_PATH.exists():
        slide2.shapes.add_picture(str(FIG1_PATH), Inches(6.9), Inches(1.9), width=Inches(5.4))

    # =========================================================================
    # SLIDE 3: THE 4-STEP PIPELINE
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "The 4-Step Chara Computational Architecture", "Methodology Overview")

    steps = [
        ("Step 1: Molecular Dynamics", "Simulate 4 oncogenic proteins across triplicate replicas using MARTINI 3 coarse-grained physics to sample allosteric ensembles.", BLUE_COLOR),
        ("Step 2: Contact Variance", "Extract residue contact variance matrices (sigma^2) and exponentially weight protein interaction graph edges.", EMERALD_COLOR),
        ("Step 3: Spectral Diffusion", "Apply the continuous graph heat diffusion operator to filter high-frequency probe noise and smooth patient profiles.", AMBER_COLOR),
        ("Step 4: Zero-Shot Survival", "Project patient expression onto the conserved 4,337-gene manifold for 5-year survival estimation and risk scoring.", NAVY_COLOR)
    ]

    for i, (title, desc, col) in enumerate(steps):
        left_pos = Inches(0.8 + i * 2.95)
        card_s = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(2.8), Inches(4.9))
        card_s.fill.solid()
        card_s.fill.fore_color.rgb = WHITE_COLOR
        card_s.line.color.rgb = col
        card_s.line.width = Pt(2)

        tx_s = slide3.shapes.add_textbox(left_pos + Inches(0.15), Inches(2.0), Inches(2.5), Inches(4.4))
        tf_s = tx_s.text_frame
        tf_s.word_wrap = True
        
        p = tf_s.paragraphs[0]
        p.text = f"0{i+1}"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = col

        p = tf_s.add_paragraph()
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY_COLOR
        p.space_before = Pt(6)

        p = tf_s.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY_TEXT
        p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 4: STEP 1 - MARTINI 3 MD SIMULATIONS
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Step 1: MARTINI 3 Coarse-Grained Molecular Dynamics", "Biophysical Simulation Setup")

    targets = [
        ("KRAS G12D (4OBE)", KRAS_PATH, "Switch I/II loop opening & cryptic pocket exploration"),
        ("c-MYC / MAX (1NKP)", CMYC_PATH, "bHLH-LZ heterodimer interface rigidity & flanking loops"),
        ("PTPN11 / SHP2 (4DGP)", PTPN11_PATH, "N-SH2 domain autoinhibition & catalytic cleft motion"),
        ("Mutant TP53 (2J1X)", P53_PATH, "DNA core domain destabilization & unfolding kinetics")
    ]

    for i, (name, img_path, desc) in enumerate(targets):
        left_pos = Inches(0.8 + i * 2.95)
        card_t = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.6), Inches(2.8), Inches(5.3))
        card_t.fill.solid()
        card_t.fill.fore_color.rgb = WHITE_COLOR
        card_t.line.color.rgb = BORDER_COLOR

        if img_path.exists():
            slide4.shapes.add_picture(str(img_path), left_pos + Inches(0.2), Inches(1.8), width=Inches(2.4), height=Inches(2.4))

        tx_t = slide4.shapes.add_textbox(left_pos + Inches(0.15), Inches(4.3), Inches(2.5), Inches(2.4))
        tf_t = tx_t.text_frame
        tf_t.word_wrap = True
        
        p = tf_t.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BLUE_COLOR

        p = tf_t.add_paragraph()
        p.text = desc
        p.font.size = Pt(10.5)
        p.font.color.rgb = GRAY_TEXT
        p.space_before = Pt(4)

    # =========================================================================
    # SLIDE 5: EMBEDDED MOLECULAR DYNAMICS VIDEO
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Live Molecular Dynamics Trajectory & Cryptic Pocket Opening", "Interactive Trajectory Analysis")

    # Video Card
    card_vid = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(8.2), Inches(5.3))
    card_vid.fill.solid()
    card_vid.fill.fore_color.rgb = RGBColor(11, 11, 12) # Dark Obsidian inside video frame
    card_vid.line.color.rgb = BLUE_COLOR
    card_vid.line.width = Pt(1.5)

    if VIDEO_PATH.exists():
        try:
            slide5.shapes.add_movie(str(VIDEO_PATH), Inches(0.9), Inches(1.7), Inches(8.0), Inches(5.1))
        except Exception as e:
            print(f"Movie embed notice: {e}")

    # Right Callout Card
    card_vcall = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(1.6), Inches(3.2), Inches(5.3))
    card_vcall.fill.solid()
    card_vcall.fill.fore_color.rgb = WHITE_COLOR
    card_vcall.line.color.rgb = BORDER_COLOR

    tx_vc = slide5.shapes.add_textbox(Inches(9.5), Inches(1.8), Inches(2.8), Inches(4.8))
    tf_vc = tx_vc.text_frame
    tf_vc.word_wrap = True
    
    p = tf_vc.paragraphs[0]
    p.text = "Trajectory Details:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLUE_COLOR

    v_points = [
        ("Target System: ", "KRAS G12D (PDB: 4OBE) coarse-grained MARTINI 3 representation."),
        ("Cryptic Pocket Opening: ", "Residues 60–75 (Switch I/II loops) dynamically explore an open cryptic pocket conformation."),
        ("Laser-Cyan Highlighting: ", "High residue contact variance (sigma^2 = 0.841) marks dynamic signaling pathways that govern downstream MAPK cascades.")
    ]
    for pref, text in v_points:
        p = tf_vc.add_paragraph()
        p.text = pref + text
        p.font.size = Pt(10.5)
        p.font.color.rgb = NAVY_COLOR
        p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 6: MD TRAJECTORY VERIFICATION (RMSD & RG)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Triplicate Trajectory Verification: RMSD & Radius of Gyration", "Step 1 Trajectory Quality")

    # RMSD Plot
    card_rmsd = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_rmsd.fill.solid()
    card_rmsd.fill.fore_color.rgb = WHITE_COLOR
    card_rmsd.line.color.rgb = BORDER_COLOR

    if RMSD_PATH.exists():
        slide6.shapes.add_picture(str(RMSD_PATH), Inches(1.0), Inches(1.9), width=Inches(5.3))

    # Rg Plot
    card_rg = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_rg.fill.solid()
    card_rg.fill.fore_color.rgb = WHITE_COLOR
    card_rg.line.color.rgb = BORDER_COLOR

    if RG_PATH.exists():
        slide6.shapes.add_picture(str(RG_PATH), Inches(7.0), Inches(1.9), width=Inches(5.3))

    # =========================================================================
    # SLIDE 7: STEP 2 - RESIDUE VARIANCE TO GRAPH
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Step 2: Residue Contact Variance to Edge Weighting", "Physical Graph Construction")

    card_g7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.3))
    card_g7.fill.solid()
    card_g7.fill.fore_color.rgb = WHITE_COLOR
    card_g7.line.color.rgb = BORDER_COLOR

    tx7 = slide7.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(10.9), Inches(4.6))
    tf7 = tx7.text_frame
    tf7.word_wrap = True

    p = tf7.paragraphs[0]
    p.text = "How Molecular Dynamics Weights Protein Networks:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_COLOR

    g_points = [
        ("1. Inter-Residue Distance Variance (sigma^2): ", "Across multi-nanosecond frames, inter-bead distance fluctuations are measured to quantify allosteric flexibility."),
        ("2. Standardized Z-Score: ", "Residue variance is standardized across the structural ensemble to isolate significant dynamic conformational hot-spots."),
        ("3. Exponential Graph Weighting: ", "Edges in the STRING protein network connecting flexible oncoprotein partners receive exponential weight amplification: W_Chara = W_STRING * exp(tau * Z(sigma^2))."),
        ("4. Allosteric Channel Prioritization: ", "Unlike unweighted graphs, signals flow preferentially through dynamically coupling allosteric switches.")
    ]
    for pref, text in g_points:
        p = tf7.add_paragraph()
        p.text = pref + text
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY_COLOR
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 8: STEP 3 - SPECTRAL HEAT DIFFUSION
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Step 3: Spectral Heat Diffusion on Biological Networks", "Manifold Filtering")

    card_d8 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.3))
    card_d8.fill.solid()
    card_d8.fill.fore_color.rgb = WHITE_COLOR
    card_d8.line.color.rgb = BORDER_COLOR

    tx8 = slide8.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(10.9), Inches(4.6))
    tf8 = tx8.text_frame
    tf8.word_wrap = True

    p = tf8.paragraphs[0]
    p.text = "Continuous Spectral Graph Filtering Principles:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_COLOR

    d_points = [
        ("Low-Pass Graph Filtering: ", "The matrix exponential H_t = exp(-t * L_sym) acts as a low-pass filter over the biological network."),
        ("Exponential Noise Attenuation: ", "High-frequency probe artifacts occupy high eigenvalues and decay exponentially: exp(-t * lambda_k) -> 0."),
        ("Pathway Smoothing: ", "True co-expression patterns among interacting pathway partners are harmonized, eliminating technical microarray spikes."),
        ("Zero Harmonization Requirement: ", "Allows unaligned single-patient transcriptomes to be processed zero-shot without cohort-level ComBat batch adjustment.")
    ]
    for pref, text in d_points:
        p = tf8.add_paragraph()
        p.text = pref + text
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY_COLOR
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 9: DIRICHLET ENERGY REDUCTION (FIGURE 5)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Empirical Proof: Dirichlet Energy Reduction in TCGA (p < 10⁻¹⁵)", "Manifold Smoothness Proof")

    card_l9 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_l9.fill.solid()
    card_l9.fill.fore_color.rgb = WHITE_COLOR
    card_l9.line.color.rgb = BORDER_COLOR

    tx9 = slide9.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.3), Inches(4.6))
    tf9 = tx9.text_frame
    tf9.word_wrap = True
    
    p = tf9.paragraphs[0]
    p.text = "Mathematical Dirichlet Smoothness:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = EMERALD_COLOR

    dir_points = [
        ("Manifold Variation: ", "Dirichlet energy E(f) = f^T * L * f measures the total variation and roughness of gene expression over the network graph."),
        ("Statistically Significant Reduction: ", "Across n=503 TCGA-LUAD patients, Chara significantly lowers topological roughness compared to standard STRING (p = 2.41 * 10^-19, paired Wilcoxon test)."),
        ("Noise Elimination: ", "Proves that incorporating atomistic MD variances actively stabilizes expression profiles against platform artifacts.")
    ]
    for pref, text in dir_points:
        p = tf9.add_paragraph()
        p.text = pref + text
        p.font.size = Pt(11.5)
        p.font.color.rgb = NAVY_COLOR
        p.space_before = Pt(8)

    card_r9 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_r9.fill.solid()
    card_r9.fill.fore_color.rgb = WHITE_COLOR
    card_r9.line.color.rgb = BORDER_COLOR

    if FIG5_PATH.exists():
        slide9.shapes.add_picture(str(FIG5_PATH), Inches(7.0), Inches(1.9), width=Inches(5.3))

    # =========================================================================
    # SLIDE 10: STEP 4 - 4,337-GENE SIGNATURE
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "Step 4: 4,337-Gene Conserved Signature & 58 Biomarkers", "Feature Space Isolation")

    gene_boxes = [
        ("19,260", "TCGA RNA-seq Transcriptome", "Full unconstrained human protein-coding transcripts.", GRAY_TEXT),
        ("12,488", "STRING Physical Graph", "Curated high-confidence protein interaction network.", GRAY_TEXT),
        ("4,337", "Conserved Manifold", "Exact mathematical intersection across RNA-seq and Affymetrix Microarrays.", BLUE_COLOR),
        ("58", "Active Biomarkers", "Non-zero prognostic genes isolated by penalized ElasticNet optimization.", EMERALD_COLOR)
    ]

    for i, (num, title, desc, col) in enumerate(gene_boxes):
        left_pos = Inches(0.8 + i * 2.95)
        card_gb = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(2.8), Inches(4.9))
        card_gb.fill.solid()
        card_gb.fill.fore_color.rgb = WHITE_COLOR
        card_gb.line.color.rgb = col
        card_gb.line.width = Pt(2)

        tx_gb = slide10.shapes.add_textbox(left_pos + Inches(0.15), Inches(2.1), Inches(2.5), Inches(4.3))
        tf_gb = tx_gb.text_frame
        tf_gb.word_wrap = True
        
        p = tf_gb.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = col

        p = tf_gb.add_paragraph()
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY_COLOR
        p.space_before = Pt(6)

        p = tf_gb.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY_TEXT
        p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 11: BENCHMARK 1 - ZERO-SHOT OOD VALIDATION (FIGURE 1)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "Benchmark 1: Zero-Shot OOD Validation on GSE31210 (n=226)", "Primary Cross-Platform Validation")

    card_l11 = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_l11.fill.solid()
    card_l11.fill.fore_color.rgb = WHITE_COLOR
    card_l11.line.color.rgb = BORDER_COLOR

    tx11 = slide11.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.3), Inches(4.6))
    tf11 = tx11.text_frame
    tf11.word_wrap = True
    
    p = tf11.paragraphs[0]
    p.text = "Held-Out Cohort Benchmarking:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = BLUE_COLOR

    table_data = [
        ("Random Survival Forest (RSF): ", "C = 0.4041 (Inverted ranking failure)"),
        ("Clinical Baseline (Age, Stage): ", "C = 0.5000 (Random coin flip)"),
        ("Standard ElasticNet: ", "C = 0.5248 (Overfitting to RNA-seq)"),
        ("DeepSurv (Deep Neural Net): ", "C = 0.5537 (Near-random performance)"),
        ("Chara Framework (Ours): ", "C = 0.7311 (Gold-standard discrimination)")
    ]
    for pref, text in table_data:
        p = tf11.add_paragraph()
        p.text = "• " + pref + text
        p.font.size = Pt(11)
        p.font.color.rgb = EMERALD_COLOR if "Chara" in pref else NAVY_COLOR
        p.font.bold = True if "Chara" in pref else False
        p.space_before = Pt(6)

    card_r11 = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_r11.fill.solid()
    card_r11.fill.fore_color.rgb = WHITE_COLOR
    card_r11.line.color.rgb = BORDER_COLOR

    if FIG1_PATH.exists():
        slide11.shapes.add_picture(str(FIG1_PATH), Inches(7.0), Inches(1.9), width=Inches(5.3))

    # =========================================================================
    # SLIDE 12: BENCHMARK 2 - TIME HORIZONS & KM SURVIVAL (FIGURE 2)
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_header(slide12, "Benchmark 2: Monotonic Time-Horizon AUCs & Kaplan-Meier Strata", "Calibration & Survival Separation")

    card_l12 = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_l12.fill.solid()
    card_l12.fill.fore_color.rgb = WHITE_COLOR
    card_l12.line.color.rgb = BORDER_COLOR

    tx12 = slide12.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.3), Inches(4.6))
    tf12 = tx12.text_frame
    tf12.word_wrap = True
    
    p = tf12.paragraphs[0]
    p.text = "Monotonic Horizon Scaling:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = BLUE_COLOR

    p = tf12.add_paragraph()
    p.text = "• 1-Year AUC: 0.7463\n• 3-Year AUC: 0.7826\n• 5-Year AUC: 0.8195"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = EMERALD_COLOR
    p.space_before = Pt(6)

    p = tf12.add_paragraph()
    p.text = "Kaplan-Meier Risk Stratification:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY_COLOR
    p.space_before = Pt(10)

    p = tf12.add_paragraph()
    p.text = "Clear statistical divergence between Low and High risk patient cohorts (log-rank p < 10^-6). Low-risk patients demonstrate 61.4% 5-year survival vs < 18.5% in high-risk groups."
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_TEXT
    p.space_before = Pt(4)

    card_r12 = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_r12.fill.solid()
    card_r12.fill.fore_color.rgb = WHITE_COLOR
    card_r12.line.color.rgb = BORDER_COLOR

    if FIG2_PATH.exists():
        slide12.shapes.add_picture(str(FIG2_PATH), Inches(7.0), Inches(1.9), width=Inches(5.3))

    # =========================================================================
    # SLIDE 13: BENCHMARK 3 - ABLATION & ADVERSARIAL STRESS (FIGURE 3 & 4)
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13)
    add_header(slide13, "Benchmark 3: Biophysical Ablation & Adversarial Stress Tests", "Robustness & Physics Impact")

    card_l13 = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_l13.fill.solid()
    card_l13.fill.fore_color.rgb = WHITE_COLOR
    card_l13.line.color.rgb = BORDER_COLOR

    if FIG3_PATH.exists():
        slide13.shapes.add_picture(str(FIG3_PATH), Inches(1.0), Inches(1.9), width=Inches(5.3))

    card_r13 = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_r13.fill.solid()
    card_r13.fill.fore_color.rgb = WHITE_COLOR
    card_r13.line.color.rgb = BORDER_COLOR

    if FIG4_PATH.exists():
        slide13.shapes.add_picture(str(FIG4_PATH), Inches(7.0), Inches(1.9), width=Inches(5.3))

    # =========================================================================
    # SLIDE 14: BIOLOGICAL ENRICHMENT (FIGURE 6) & BIOMARKERS
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14)
    add_header(slide14, "Biological GSEA Hallmarks & Key Prognostic Drivers", "Mechanistic Validation")

    card_l14 = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_l14.fill.solid()
    card_l14.fill.fore_color.rgb = WHITE_COLOR
    card_l14.line.color.rgb = BORDER_COLOR

    if FIG6_PATH.exists():
        slide14.shapes.add_picture(str(FIG6_PATH), Inches(1.0), Inches(1.9), width=Inches(5.3))

    card_r14 = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_r14.fill.solid()
    card_r14.fill.fore_color.rgb = WHITE_COLOR
    card_r14.line.color.rgb = BORDER_COLOR

    tx14_r = slide14.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.3), Inches(4.6))
    tf14_r = tx14_r.text_frame
    tf14_r.word_wrap = True

    p = tf14_r.paragraphs[0]
    p.text = "Key Active Prognostic Biomarkers:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY_COLOR

    b_drivers = [
        ("CCL20 (beta = +0.0642): ", "Recruits CCR6+ regulatory T-cells to promote immune evasion."),
        ("DKK1 (beta = +0.0610): ", "Wnt signaling inhibitor linked to osteolytic bone metastasis."),
        ("MS4A1 / CD20 (beta = -0.0708): ", "B-lymphocyte antigen indicating active tertiary lymphoid structure immunity."),
        ("FAIM2 (beta = -0.0524): ", "Homeostatic regulator associated with well-differentiated phenotype.")
    ]
    for pref, text in b_drivers:
        p = tf14_r.add_paragraph()
        p.text = "• " + pref + text
        p.font.size = Pt(11)
        p.font.color.rgb = CRIMSON_COLOR if "beta = +" in pref else EMERALD_COLOR
        p.space_before = Pt(6)

    # =========================================================================
    # SLIDE 15: CLINICAL STRATIFICATION & SOFTWARE ECOSYSTEM
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15)
    add_header(slide15, "Clinical Risk Stratification & Open-Source Ecosystem", "Translation & Software")

    # Left Column Risk Table
    card_l15 = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_l15.fill.solid()
    card_l15.fill.fore_color.rgb = WHITE_COLOR
    card_l15.line.color.rgb = BORDER_COLOR

    tx15_l = slide15.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.3), Inches(4.6))
    tf15_l = tx15_l.text_frame
    tf15_l.word_wrap = True

    p = tf15_l.paragraphs[0]
    p.text = "4-Tier Clinical Stratification:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = BLUE_COLOR

    r_tiers = [
        ("Low Risk (Score < -0.50): ", "5-Yr Survival: ~61.4%. Surveillance candidate."),
        ("Moderate Risk (-0.50 to +0.50): ", "5-Yr Survival: ~35.2%. Standard clinical baseline."),
        ("High Risk (+0.50 to +1.20): ", "5-Yr Survival: ~18.5%. Recommended for targeted adjuvant therapy."),
        ("Critical Risk (Score > +1.20): ", "5-Yr Survival: < 8.2%. Candidate for aggressive trial regimens.")
    ]
    for pref, text in r_tiers:
        p = tf15_l.add_paragraph()
        p.text = "• " + pref + text
        p.font.size = Pt(10.5)
        p.font.color.rgb = NAVY_COLOR
        p.space_before = Pt(6)

    p = tf15_l.add_paragraph()
    p.text = "Multivariate Independence: HR = 3.81 (p = 4.48 * 10^-33) adjusted for Age, Gender, and TNM Stage."
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD_COLOR
    p.space_before = Pt(8)

    # Right Column Ecosystem
    card_r15 = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.3))
    card_r15.fill.solid()
    card_r15.fill.fore_color.rgb = WHITE_COLOR
    card_r15.line.color.rgb = BORDER_COLOR

    tx15_r = slide15.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.3), Inches(4.6))
    tf15_r = tx15_r.text_frame
    tf15_r.word_wrap = True

    p = tf15_r.paragraphs[0]
    p.text = "Production Software Ecosystem:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = EMERALD_COLOR

    eco_items = [
        ("Official PyPI Library (v0.2.5): ", "pip install chara-survival (over 4,600+ downloads)"),
        ("Hugging Face Hub: ", "SharonMelhi/chara-survival pretrained weights (chara_model_4337.pkl)"),
        ("Live Web Application: ", "https://chara-frontend.vercel.app (Client-side WebAssembly UI)"),
        ("Institutional Affiliation: ", "Computational & Physical Genomics Lab, IIT Mandi")
    ]
    for pref, text in eco_items:
        p = tf15_r.add_paragraph()
        p.text = "• " + pref + text
        p.font.size = Pt(10.5)
        p.font.color.rgb = NAVY_COLOR
        p.space_before = Pt(8)

    # Save outputs
    output_desktop = Path(r"C:\Users\Samsunh\Desktop\Chara_Symposium_Presentation.pptx")
    output_workspace = ROOT / "Chara_Symposium_Presentation.pptx"

    prs.save(str(output_desktop))
    prs.save(str(output_workspace))
    print(f"SUCCESS: Generated Master PowerPoint presentation at:\n -> {output_desktop}\n -> {output_workspace}")

if __name__ == "__main__":
    create_presentation()
